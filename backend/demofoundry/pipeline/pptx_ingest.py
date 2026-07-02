"""PowerPoint ingestion — turn a .pptx into a narratable DemoFoundry deck.

Pipeline:
  1. Export each slide to a faithful 1920x1080 PNG via PowerPoint COM automation
     (requires PowerPoint installed — Windows + Microsoft Office).
  2. Read each slide's speaker notes / text via python-pptx.
  3. Narrate: speaker notes win; otherwise Claude writes spoken narration from the
     slide's text (skipped with narrate=False).
  4. Emit a minimal full-bleed HTML image-deck with a forward button (#btn-next)
     plus a steps.json (navigate + one click per advance). The deck reuses the
     existing Playwright capture + compose path unchanged — one scene per slide,
     the same step schema as a web demo, with narration_text the only meaningful
     per-slide field.

Cross-platform note: this ingester uses two backends:
  - PowerPoint COM on Windows, when available.
  - LibreOffice headless (`soffice`) -> PDF -> PNG rasterization for macOS/Linux.
`powerpoint_available()` lets callers gate the feature when neither backend is
present.
"""

from __future__ import annotations

import json
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

from .. import config

# Slides export at the same 16:9 frame DemoFoundry composes to.
SLIDE_W, SLIDE_H = 1920, 1080


def _windows_powerpoint_available() -> bool:
    """True if PowerPoint COM automation is usable on this machine.

    Lightweight: checks the COM registration in the registry rather than
    launching PowerPoint. Always False off-Windows.
    """
    if sys.platform != "win32":
        return False
    try:
        import winreg
    except ImportError:
        return False
    try:
        import win32com.client  # noqa: F401  (presence check)
    except ImportError:
        return False
    try:
        winreg.CloseKey(winreg.OpenKey(winreg.HKEY_CLASSES_ROOT, "PowerPoint.Application"))
        return True
    except OSError:
        return False


def _libreoffice_available() -> bool:
    """True if LibreOffice headless conversion + PDF rasterizer are available."""
    if not shutil.which("soffice"):
        return False
    try:
        import pypdfium2  # noqa: F401  (presence check)
    except ImportError:
        return False
    return True


def powerpoint_available() -> bool:
    """True if any PPTX slide-export backend is usable on this machine."""
    return _windows_powerpoint_available() or _libreoffice_available()


def _export_slides_windows(pptx_path: Path, out_dir: Path) -> list[Path]:
    """Export every slide to out_dir/slide-NN.png at 1920x1080 via PowerPoint COM.

    Safe against a PowerPoint the user already has open: we only close the
    presentation we opened, and only quit the app if we were the ones to start it.
    """
    import pythoncom
    import win32com.client

    pptx_path = Path(pptx_path).resolve()
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    for old in out_dir.glob("slide-*.png"):
        old.unlink()

    pythoncom.CoInitialize()
    app = win32com.client.Dispatch("PowerPoint.Application")
    started_clean = False
    paths: list[Path] = []
    try:
        # If no presentations are open, this process effectively launched PP and
        # is responsible for quitting it; otherwise leave the user's session be.
        started_clean = app.Presentations.Count == 0
        pres = app.Presentations.Open(str(pptx_path), WithWindow=False)
        try:
            for i, slide in enumerate(pres.Slides, 1):
                dest = out_dir / f"slide-{i:02d}.png"
                slide.Export(str(dest), "PNG", SLIDE_W, SLIDE_H)
                paths.append(dest)
        finally:
            pres.Close()
    finally:
        if started_clean:
            app.Quit()
        pythoncom.CoUninitialize()
    return paths


def _export_slides_libreoffice(pptx_path: Path, out_dir: Path) -> list[Path]:
    """Export slides via LibreOffice headless conversion to PDF + rasterization."""
    if not _libreoffice_available():
        raise RuntimeError(
            "LibreOffice path unavailable: install LibreOffice (soffice on PATH) "
            "and pypdfium2."
        )
    import pypdfium2 as pdfium

    pptx_path = Path(pptx_path).resolve()
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    for old in out_dir.glob("slide-*.png"):
        old.unlink()

    with tempfile.TemporaryDirectory(prefix="demofoundry-pptx-") as td:
        td_path = Path(td)
        cmd = [
            "soffice",
            "--headless",
            "--nologo",
            "--nodefault",
            "--nofirststartwizard",
            "--convert-to",
            "pdf:impress_pdf_Export",
            "--outdir",
            str(td_path),
            str(pptx_path),
        ]
        try:
            subprocess.run(cmd, check=True, capture_output=True, text=True)
        except subprocess.CalledProcessError as e:
            detail = (e.stderr or e.stdout or "").strip()
            raise RuntimeError(f"LibreOffice conversion failed: {detail or e}") from e

        pdf_path = td_path / f"{pptx_path.stem}.pdf"
        if not pdf_path.exists():
            pdfs = sorted(td_path.glob("*.pdf"))
            if not pdfs:
                raise RuntimeError("LibreOffice did not produce a PDF from the PPTX.")
            pdf_path = pdfs[0]

        doc = pdfium.PdfDocument(str(pdf_path))
        paths: list[Path] = []
        from PIL import Image

        for i in range(len(doc)):
            page = doc[i]
            try:
                width_pt, _ = page.get_size()
                scale = SLIDE_W / width_pt if width_pt > 0 else 1.0
                img = page.render(scale=scale).to_pil().convert("RGB")
            finally:
                page.close()
            if img.size != (SLIDE_W, SLIDE_H):
                fit = min(SLIDE_W / img.width, SLIDE_H / img.height)
                size = (max(1, int(img.width * fit)), max(1, int(img.height * fit)))
                fg = img.resize(size)
                bg = Image.new("RGB", (SLIDE_W, SLIDE_H), (0, 0, 0))
                bg.paste(fg, ((SLIDE_W - size[0]) // 2, (SLIDE_H - size[1]) // 2))
                img = bg
            dest = out_dir / f"slide-{i + 1:02d}.png"
            img.save(dest, format="PNG")
            paths.append(dest)
        doc.close()
        if not paths:
            raise RuntimeError("No slide images were generated from the converted PDF.")
        return paths


def export_slides(pptx_path: Path, out_dir: Path) -> list[Path]:
    """Export every slide to out_dir/slide-NN.png at 1920x1080."""
    if _windows_powerpoint_available():
        return _export_slides_windows(pptx_path, out_dir)
    if _libreoffice_available():
        return _export_slides_libreoffice(pptx_path, out_dir)
    raise RuntimeError(
        "No PPTX export backend available. Install Microsoft PowerPoint (Windows), "
        "or install LibreOffice (soffice) plus pypdfium2."
    )


def extract_slides(pptx_path: Path) -> list[dict]:
    """Per-slide {title, body, notes} pulled from the .pptx via python-pptx.

    `title` is the slide's title placeholder if present; `body` is the remaining
    shape text joined; `notes` is the speaker-notes text (often empty).
    """
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    slides: list[dict] = []
    for slide in prs.slides:
        title = ""
        if slide.shapes.title is not None:
            title = (slide.shapes.title.text or "").strip()
        body_parts: list[str] = []
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.has_text_frame:
                txt = (shape.text_frame.text or "").strip()
                if txt:
                    body_parts.append(txt)
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        slides.append({"title": title, "body": "\n".join(body_parts), "notes": notes})
    return slides


_NARRATE_SYSTEM = (
    "You write spoken voiceover narration for a slide presentation video. For each "
    "slide you get its on-screen text. Write two or three natural, spoken sentences "
    "per slide — what the presenter would say while that slide is up. Carry the "
    "thread from one slide to the next so it plays as one continuous talk. Don't "
    "read bullets verbatim or say 'this slide'; speak to the point. No stage "
    "directions."
)

_NARRATE_SCHEMA = {
    "type": "object",
    "properties": {
        "narration": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "slide": {"type": "integer"},
                    "narration_text": {"type": "string"},
                },
                "required": ["slide", "narration_text"],
                "additionalProperties": False,
            },
        }
    },
    "required": ["narration"],
    "additionalProperties": False,
}


def narrate_slides(slides: list[dict], deck_title: str = "") -> list[str]:
    """Claude writes spoken narration for slides lacking speaker notes.

    Returns one narration string per slide (notes are used verbatim where present;
    only the blank ones are sent to Claude). Requires ANTHROPIC_API_KEY.
    """
    out = [s.get("notes", "").strip() for s in slides]
    blank_idx = [i for i, n in enumerate(out) if not n]
    if not blank_idx:
        return out
    if not config.ANTHROPIC_API_KEY:
        raise RuntimeError("ANTHROPIC_API_KEY not set — needed to narrate slides without speaker notes")

    import anthropic

    brief = [
        {"slide": i + 1, "title": slides[i]["title"], "text": slides[i]["body"]}
        for i in blank_idx
    ]
    user = (
        f"Presentation: {deck_title or '(untitled)'}\n\n"
        f"Write narration ONLY for these slides (by their slide number):\n"
        f"{json.dumps(brief, indent=2)}\n\n"
        "Return one entry per requested slide, using its slide number."
    )
    client = anthropic.Anthropic(api_key=config.ANTHROPIC_API_KEY)
    resp = client.messages.create(
        model=config.CLAUDE_MODEL,
        max_tokens=8000,
        thinking={"type": "adaptive"},
        system=_NARRATE_SYSTEM,
        messages=[{"role": "user", "content": user}],
        output_config={"format": {"type": "json_schema", "schema": _NARRATE_SCHEMA}},
    )
    text = next(b.text for b in resp.content if b.type == "text")
    by_slide = {e["slide"]: e["narration_text"] for e in json.loads(text).get("narration", [])}
    for i in blank_idx:
        out[i] = by_slide.get(i + 1, "").strip()
    return out


_DECK_TEMPLATE = """<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8">
<title>__TITLE__</title>
<style>
  * { margin: 0; padding: 0; box-sizing: border-box; }
  html, body { width: 100%; height: 100%; background: #000; overflow: hidden; }
  #stage { position: fixed; inset: 0; display: flex; align-items: center; justify-content: center; }
  .slide { position: absolute; max-width: 100%; max-height: 100%; opacity: 0;
           transition: opacity 0.25s ease; pointer-events: none; }
  .slide.active { opacity: 1; }
  #nav { position: fixed; bottom: 14px; left: 50%; transform: translateX(-50%);
         display: flex; gap: 10px; align-items: center; z-index: 10;
         font: 13px system-ui, sans-serif; }
  #nav button { background: #1f6feb; color: #fff; border: none; border-radius: 6px;
                padding: 7px 14px; cursor: pointer; font-weight: 600; }
  #nav button:disabled { background: #30363d; color: #8b949e; cursor: default; }
  #counter { color: #c9d1d9; min-width: 96px; text-align: center; }
</style>
</head>
<body>
  <div id="stage">__SLIDES__</div>
  <div id="nav">
    <button id="btn-back" title="Previous (Left Arrow)">‹ Back</button>
    <span id="counter"></span>
    <button id="btn-next" title="Next (Right Arrow)">Next ›</button>
  </div>
<script>
  const slides = Array.from(document.querySelectorAll('.slide'));
  const total = slides.length;
  let i = 0;
  const counter = document.getElementById('counter');
  const back = document.getElementById('btn-back');
  const next = document.getElementById('btn-next');
  function render() {
    slides.forEach((s, n) => s.classList.toggle('active', n === i));
    counter.textContent = `Slide ${i + 1} of ${total}`;
    back.disabled = i === 0;
    next.disabled = i === total - 1;
  }
  function go(n) { i = Math.max(0, Math.min(total - 1, n)); render(); }
  next.addEventListener('click', () => go(i + 1));
  back.addEventListener('click', () => go(i - 1));
  // During capture the deck is loaded with #capture so the nav chrome is hidden
  // and slides render clean; advancing is driven by ArrowRight keypresses.
  if (location.hash.includes('capture')) document.getElementById('nav').style.display = 'none';
  window.addEventListener('keydown', (e) => {
    if (e.key === 'ArrowRight' || e.key === 'ArrowDown') go(i + 1);
    else if (e.key === 'ArrowLeft' || e.key === 'ArrowUp') go(i - 1);
    else if (e.key === 'Home') go(0);
    else if (e.key === 'End') go(total - 1);
  });
  render();
</script>
</body>
</html>
"""


def build_deck_html(png_paths: list[Path], out_path: Path, title: str = "Slides") -> Path:
    """Write a minimal full-bleed image slideshow (forward button = #btn-next).

    Images are referenced relatively, so keep the PNGs beside the HTML (the
    default layout from ingest()).
    """
    out_path = Path(out_path)
    # Reference each PNG by path relative to the HTML file.
    rels = [Path(p).resolve().relative_to(out_path.parent.resolve()).as_posix()
            if Path(p).resolve().is_relative_to(out_path.parent.resolve())
            else Path(p).name for p in png_paths]
    imgs = "\n    ".join(
        f'<img class="slide{" active" if n == 0 else ""}" src="{rel}" alt="Slide {n + 1}">'
        for n, rel in enumerate(rels)
    )
    html = _DECK_TEMPLATE.replace("__TITLE__", title).replace("__SLIDES__", imgs)
    out_path.write_text(html, encoding="utf-8")
    return out_path


def build_steps(deck_url: str, narration: list[str], name: str) -> dict:
    """steps.json for the image-deck: navigate to slide 1 (with #capture so the
    nav chrome is hidden), then advance one slide per ArrowRight keypress. Keypress
    leaves no click marker, so slides render clean. Same schema as a web demo."""
    capture_url = deck_url + "#capture"
    steps = [{
        "action": "navigate",
        "value": capture_url,
        "narration_text": narration[0] if narration else "",
    }]
    for n in narration[1:]:
        steps.append({"action": "keypress", "value": "ArrowRight", "narration_text": n})
    return {"name": name, "target_url": capture_url, "steps": steps}


def ingest(pptx_path: Path, out_dir: Path, *, narrate: bool = True) -> dict:
    """Full ingest: export slides, extract text/notes, narrate, write deck.html +
    steps.json into out_dir. Returns a summary dict with the output paths."""
    pptx_path = Path(pptx_path)
    out_dir = Path(out_dir)
    slides_dir = out_dir / "slides"
    pngs = export_slides(pptx_path, slides_dir)

    meta = extract_slides(pptx_path)
    # extract_slides and export must agree on count; trust the export count.
    while len(meta) < len(pngs):
        meta.append({"title": "", "body": "", "notes": ""})
    meta = meta[: len(pngs)]

    title = pptx_path.stem
    narration = (
        narrate_slides(meta, deck_title=title) if narrate
        else [m.get("notes", "").strip() for m in meta]
    )

    deck_html = build_deck_html(pngs, out_dir / "deck.html", title=title)
    deck_url = deck_html.resolve().as_uri()
    steps = build_steps(deck_url, narration, name=title)
    steps_path = out_dir / "steps.json"
    steps_path.write_text(json.dumps(steps, indent=2), encoding="utf-8")

    return {
        "slides": len(pngs),
        "deck_html": str(deck_html),
        "deck_url": deck_url,
        "steps": str(steps_path),
        "narrated": narrate,
        "notes_found": sum(1 for m in meta if m.get("notes", "").strip()),
    }
